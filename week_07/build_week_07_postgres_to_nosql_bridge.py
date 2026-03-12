from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Week_07_Postgres_to_NoSQL_Bridge.pptx"

NAVY = RGBColor(9, 21, 39)
NAVY_LIGHT = RGBColor(17, 45, 78)
NAVY_PANEL = RGBColor(21, 55, 95)
GOLD = RGBColor(247, 191, 86)
TEAL = RGBColor(123, 224, 214)
WHITE = RGBColor(245, 247, 250)
MUTED = RGBColor(199, 210, 229)
ROSE = RGBColor(255, 204, 182)


SLIDES = [
    {
        "kind": "title",
        "tag": "Week 7 Bridge",
        "title": "From PostgreSQL to NoSQL",
        "subtitle": (
            "Use the last part of Week 7 to explain why NoSQL re-emerged, what the main "
            "NoSQL families are, and why Week 8 will narrow to MongoDB."
        ),
        "focus_title": "Bridge goals",
        "focus_lines": [
            "Explain the historical shift",
            "Classify the major NoSQL families",
            "Set up MongoDB as next week's document example",
        ],
        "footer": "This short deck is the transition from the relational unit into the MongoDB unit.",
        "notes": """Use this as the closing bridge after the PostgreSQL wrap-up, not as a full MongoDB lecture. The goal is to widen the students' mental model before Week 8, not to start a new syntax-heavy unit too early.

Keep the contrast clean: Weeks 1 through 7 taught relational design, PostgreSQL operations, concurrency, and maintenance. The bridge asks a new question: what happens when the workload does not fit the default relational shape as neatly?

Students should leave Week 7 with a map of why NoSQL became important and what the major families are. Week 8 will then spend its time on MongoDB-specific modeling and query patterns.""",
    },
    {
        "kind": "content",
        "tag": "Transition",
        "title": "Why End Week 7 With a NoSQL Bridge?",
        "bullets": [
            "We just finished a strong relational and PostgreSQL foundation",
            "Modern DBAs rarely work with only one data model in production",
            "Some workloads stress joins, rigid schemas, or single-node scaling assumptions",
            "The point is not to replace SQL, but to understand where another model fits better",
            "Week 8 will take the document-database branch of this story through MongoDB",
        ],
        "terms": ["relational", "polyglot persistence", "workload fit", "document", "DBA mindset"],
        "callout": "This is a bridge slide: the course is expanding scope, not abandoning PostgreSQL.",
        "footer": "Frame the transition as broader database literacy rather than a product switch.",
        "notes": """Students can misread this transition if it is rushed. Say plainly that PostgreSQL remains important and that the course is not declaring relational databases outdated.

What changes is the scope of the DBA mindset. A strong administrator should be able to explain when a document model, graph model, or key-value model better fits the problem at hand. That is what makes this bridge useful.

If you want a one-line summary for students, use this: 'We are not leaving databases behind; we are learning that databases come in multiple shapes for multiple workloads.'""",
    },
    {
        "kind": "content",
        "tag": "History",
        "title": "How NoSQL Re-Emerged",
        "bullets": [
            "Early databases included hierarchical and network models before relational systems dominated",
            "The relational model won because SQL and declarative querying made data easier to reason about",
            "Large web applications later stressed object-table mismatch, scaling limits, and rapid schema change",
            "Systems like Bigtable and Dynamo made distributed storage patterns central again",
            "The NoSQL movement gathered tools focused on scale, fault tolerance, and flexible records",
        ],
        "terms": ["hierarchical", "relational", "Dynamo", "Bigtable", "NoSQL"],
        "callout": "NoSQL is not brand new. It is an old set of ideas made urgent by newer workloads.",
        "footer": "History helps students see NoSQL as a response to workload pressure rather than hype.",
        "notes": """Do not over-teach the timeline here. The purpose is to correct the simplified story that everything was relational until one vendor invented something modern.

The better story is that relational systems solved major problems elegantly, but later internet-scale workloads reopened older questions under new conditions such as distribution, replication, and flexible application-shaped records.

You only need enough history for students to understand cause and effect. The detailed MongoDB content belongs in Week 8.""",
    },
    {
        "kind": "content",
        "tag": "Families",
        "title": "The Main NoSQL Families",
        "bullets": [
            "Key-value: best when you already know the key and need extremely fast retrieval",
            "Document: best when the application works with nested object-like records",
            "Column-family: best for very large sparse or time-windowed write-heavy workloads",
            "Graph: best when relationships and traversals are the query itself",
            "Choose by query shape, consistency needs, and operational constraints",
        ],
        "terms": ["key-value", "document", "column-family", "graph", "query shape"],
        "callout": "Ask 'What does the workload ask for?' before asking 'Which product should I use?'",
        "footer": "By the end of this slide, students should be able to classify a workload into a family.",
        "notes": """This slide should give students a compact territory map. Keep one workload image attached to each family so the terms do not stay abstract.

Key-value is 'return by key fast.' Document is 'store an application-shaped record.' Column-family is 'high-volume wide or time-windowed data.' Graph is 'follow relationships efficiently.'

At this point, tell students that Week 8 will focus on the document family because MongoDB is a practical and teachable next step from the course's existing database knowledge.""",
    },
    {
        "kind": "code",
        "tag": "Preview",
        "title": "Where MongoDB Fits in That Landscape",
        "bullets": [
            "MongoDB is the document-database example in this unit",
            "It stores BSON documents that can nest arrays and subdocuments",
            "It is useful when related data is usually read together as one aggregate",
            "It still needs indexes, validation, backup/restore planning, and security",
        ],
        "code_title": "Teaser document",
        "code": """{
  customerId: 42,
  status: "shipped",
  shippingCity: "Tampa",
  lineItems: [
    { sku: "SSD-1TB", qty: 1, price: 89.99 },
    { sku: "USB-C-HUB", qty: 2, price: 24.50 }
  ]
}""",
        "footer": "Week 8 will start from this document shape and compare it against a relational design.",
        "notes": """Keep this slide at teaser level. The point is to give students one concrete image of a document so Week 8 does not begin from zero.

Say clearly that a document database is not an excuse to stop modeling. The reason MongoDB is useful is that some workloads naturally read and write a whole aggregate together. That is the modeling question Week 8 will answer in more detail.

Connect back to the DBA role. Even in a flexible document system, you still have to think about validation, indexing, recovery, and observability.""",
    },
    {
        "kind": "content",
        "tag": "Next Week",
        "title": "What Week 8 Will Cover",
        "bullets": [
            "Relational thinking versus document thinking",
            "Embed versus reference decisions",
            "MongoDB CRUD and aggregation basics",
            "MongoDB trade-offs, anti-patterns, and DBA responsibilities",
            "A short review packet that spans the bridge and the MongoDB lecture",
        ],
        "terms": ["modeling", "embed", "reference", "aggregation", "review"],
        "callout": "Exit prompt: name one workload that feels more relational and one that feels more document-shaped.",
        "footer": "Use the final minute to make students predict fit, not just copy definitions.",
        "notes": """End the bridge by making students anticipate the next class. If they can already name a relational-shaped workload and a document-shaped workload, they are ready for Week 8.

You do not need to resolve every modeling question here. The goal is to create curiosity and enough vocabulary that the next lecture can move faster without feeling abrupt.

If you assign the review packet after Week 7, tell students they only need to attempt the history and classification pieces before Week 8. The MongoDB-specific sections can wait until after the next class.""",
    },
]


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()

    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(0.42), Inches(3.45), Inches(0.46)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = NAVY_LIGHT
    glow.line.color.rgb = TEAL
    glow.line.width = Pt(1.0)


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.85), Inches(0.26))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_notes(slide, notes_text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes_text.strip()


def add_tag(slide, text, width=1.75):
    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.45), Inches(width), Inches(0.38)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_box.line.width = Pt(1.0)

    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL


def add_title_text(slide, text, width=10.0):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), Inches(width), Inches(0.75))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE


def populate_bullets(text_frame, bullets, font_size=19, color=WHITE, bullet_prefix="- "):
    text_frame.word_wrap = True
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"{bullet_prefix}{bullet}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)


def add_terms_panel(slide, terms):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.78), Inches(1.58), Inches(3.62), Inches(4.9)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.3)

    tf = panel.text_frame
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = "Key terms"
    header.font.name = "Aptos"
    header.font.size = Pt(16)
    header.font.bold = True
    header.font.color.rgb = GOLD

    for term in terms:
        p = tf.add_paragraph()
        p.text = term
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(5)


def add_callout(slide, text):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(6.28), Inches(11.52), Inches(0.46)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY_LIGHT
    box.line.color.rgb = ROSE
    box.line.width = Pt(1.2)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.62), Inches(2.2), Inches(0.44)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_box.line.width = Pt(1.2)

    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = spec["tag"]
    tag_p.alignment = PP_ALIGN.CENTER
    tag_p.font.name = "Aptos"
    tag_p.font.size = Pt(14)
    tag_p.font.bold = True
    tag_p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.28), Inches(10.8), Inches(1.45))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = spec["title"]
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.62), Inches(9.8), Inches(1.42))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    p = subtitle_tf.paragraphs[0]
    p.text = spec["subtitle"]
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(3.82), Inches(3.35), Inches(1.96)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.4)

    tf = panel.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = spec["focus_title"]
    p1.font.name = "Aptos"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = GOLD

    for line in spec["focus_lines"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_content_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, spec["tag"])
    add_title_text(slide, spec["title"])

    bullet_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(7.35), Inches(4.6))
    populate_bullets(bullet_box.text_frame, spec["bullets"], font_size=18)
    add_terms_panel(slide, spec["terms"])

    if spec.get("callout"):
        add_callout(slide, spec["callout"])

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_code_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, spec["tag"], width=1.6)
    add_title_text(slide, spec["title"], width=10.7)

    bullet_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(4.55), Inches(4.72))
    populate_bullets(bullet_box.text_frame, spec["bullets"], font_size=17)

    code_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.62), Inches(1.72), Inches(6.7), Inches(4.9)
    )
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = NAVY_PANEL
    code_panel.line.color.rgb = TEAL
    code_panel.line.width = Pt(1.2)

    tf = code_panel.text_frame
    tf.word_wrap = True
    tf.clear()

    header = tf.paragraphs[0]
    header.text = spec["code_title"]
    header.font.name = "Aptos"
    header.font.size = Pt(16)
    header.font.bold = True
    header.font.color.rgb = TEAL
    header.space_after = Pt(8)

    body = tf.add_paragraph()
    body.text = spec["code"]
    body.font.name = "Courier New"
    body.font.size = Pt(13)
    body.font.color.rgb = WHITE

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        if spec["kind"] == "title":
            add_title_slide(prs, spec)
        elif spec["kind"] == "content":
            add_content_slide(prs, spec)
        elif spec["kind"] == "code":
            add_code_slide(prs, spec)
        else:
            raise ValueError(f"Unsupported slide kind: {spec['kind']}")

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Wrote {OUTPUT}")
