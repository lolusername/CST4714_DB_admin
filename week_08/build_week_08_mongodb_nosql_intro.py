from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "mongodb_nosql_intro.pptx"

NAVY = RGBColor(9, 21, 39)
NAVY_LIGHT = RGBColor(17, 45, 78)
NAVY_PANEL = RGBColor(21, 55, 95)
GOLD = RGBColor(247, 191, 86)
TEAL = RGBColor(123, 224, 214)
WHITE = RGBColor(245, 247, 250)
MUTED = RGBColor(199, 210, 229)
ROSE = RGBColor(255, 204, 182)


SLIDES = [
    {
        "kind": "title",
        "tag": "Week 8",
        "title": "MongoDB Foundations",
        "subtitle": (
            "Week 8 picks up after the Week 7 NoSQL bridge and focuses on MongoDB-specific "
            "modeling, querying, trade-offs, and administration."
        ),
        "focus_title": "Teaching focus",
        "focus_lines": [
            "Quick recap from Week 7",
            "MongoDB modeling decisions",
            "CRUD, aggregation, and admin trade-offs",
        ],
        "footer": (
            "Week 7 now carries the broad NoSQL bridge; Week 8 narrows to MongoDB itself."
        ),
        "notes": """This lecture no longer needs to carry the full NoSQL history by itself because the broad bridge now happens at the end of Week 7. That means Week 8 can spend more time on MongoDB-specific modeling and operations.

Begin with a fast recap so students reconnect the NoSQL idea to MongoDB without retaking the whole bridge. Then push into the design questions they actually need: what does a document look like, when should data be embedded, how do CRUD and aggregation feel, and what responsibilities stay with the DBA?

The tone should still be comparative and practical. MongoDB is the example, but the larger lesson is about workload fit and trade-offs.""",
    },
    {
        "kind": "content",
        "tag": "Recap",
        "title": "What Carries Over From the Week 7 Bridge",
        "bullets": [
            "NoSQL re-emerged because some web-scale and cloud workloads stressed relational defaults",
            "The four major families are key-value, document, column-family, and graph",
            "MongoDB is the document-database example for this unit",
            "We still choose databases by workload fit, not by fashion",
            "The DBA mindset carries forward even when the data model changes",
        ],
        "terms": ["NoSQL", "document", "polyglot persistence", "workload fit", "DBA"],
        "callout": "Use this as a 3-5 minute reconnect, not as a second full history lecture.",
        "footer": "The broad NoSQL bridge now happens in Week 7 so Week 8 can move faster.",
        "notes": """This slide exists to reconnect the class, not to reteach everything from the bridge deck.

The key reminders are simple: NoSQL re-emerged because some workloads made different trade-offs attractive, MongoDB is the document-database case for this unit, and relational thinking is still relevant because comparison is what makes the model choices understandable.

When students hear 'schema-less,' correct it again here. MongoDB gives you flexibility, but it still requires structure, review, and operational discipline.""",
    },
    {
        "kind": "code",
        "tag": "MongoDB",
        "title": "Where MongoDB Fits",
        "bullets": [
            "A MongoDB server holds databases, databases hold collections, and collections hold documents",
            "Documents are BSON records: JSON-like, but with richer native types such as ObjectId and Date",
            "MongoDB is strongest when the application naturally reads and writes related data together",
            "Flexible structure does not remove governance; validation, indexes, and backups still matter",
        ],
        "code_title": "Sample order document",
        "code": """{
  _id: ObjectId("65fb..."),
  customerId: 42,
  customerName: "A. Rivera",
  status: "shipped",
  shippingCity: "Tampa",
  lineItems: [
    { sku: "SSD-1TB", qty: 1, price: 89.99 },
    { sku: "USB-C-HUB", qty: 2, price: 24.50 }
  ],
  orderedAt: ISODate("2026-03-10T14:22:00Z")
}""",
        "footer": "MongoDB's document model mirrors application objects more directly than a fully normalized schema.",
        "notes": """Pause on vocabulary here because students new to MongoDB often lose confidence when the terms pile up. Keep the stack simple: server, database, collection, document, field.

The sample document is important because it shows why document databases feel intuitive to many developers. The order screen in an application often needs the customer snapshot, status, shipping information, and line items at the same time. A document can package that read pattern neatly.

At the same time, warn students not to treat a document as a junk drawer. The fact that MongoDB permits nested structure does not mean every related fact should be embedded forever. Modeling discipline still matters.""",
    },
    {
        "kind": "compare",
        "tag": "Modeling",
        "title": "Relational Thinking vs Document Thinking",
        "left_title": "Relational approach",
        "left_bullets": [
            "Separate `customers`, `orders`, and `order_items` tables",
            "Rebuild the order view with joins when the app needs a complete screen",
            "Minimize duplication and enforce shared facts centrally",
            "Best when related facts change independently or support many cross-entity queries",
        ],
        "right_title": "Document approach",
        "right_bullets": [
            "Store an `orders` document with embedded line items and shipping snapshot",
            "Read the whole order screen with fewer joins and fewer round trips",
            "Accept intentional duplication if it improves the dominant read path",
            "Best when the aggregate is usually created, read, and updated together",
        ],
        "callout": "The design question is not 'Which is modern?' It is 'Which shape matches the workload?'",
        "footer": "Same business domain, different storage shape because the dominant queries differ.",
        "notes": """This is one of the most important teaching slides in the deck. Students already know normalization, so show them that document modeling is not random. It has a logic. The logic is to design around aggregates and access patterns rather than decomposing everything into reusable tables by default.

Use the order example because it makes duplication acceptable in a way students can reason about. Shipping addresses and product names often need to be preserved as historical snapshots anyway. That makes the document model easier to justify.

Also show the limit of that logic. If the same customer record must be updated centrally across many workflows, or if the workload involves many cross-order analytical joins, a relational structure may still be better.""",
    },
    {
        "kind": "compare",
        "tag": "Design",
        "title": "Embed or Reference?",
        "left_title": "Embed when...",
        "left_bullets": [
            "The child data is read with the parent almost every time",
            "The relationship is one-to-few and size stays bounded",
            "Updates happen as one business unit",
            "Duplication is acceptable or even desirable for historical snapshots",
        ],
        "right_title": "Reference when...",
        "right_bullets": [
            "The child data has an independent lifecycle or many parents",
            "Arrays can grow without a safe bound",
            "The workload needs separate updates, ownership, or sharing",
            "Many-to-many relationships would create large duplicated structures",
        ],
        "callout": (
            "Quick rule: bounded and read-together leans embed; shared or unbounded leans reference."
        ),
        "footer": "Good MongoDB design starts with honest read/write patterns, not with maximum nesting.",
        "notes": """Students usually overcorrect in one of two directions: they either try to normalize MongoDB until it looks like tables, or they embed everything and create giant documents.

Teach the embed/reference choice as a boundedness and ownership question. If the child data belongs to the parent and stays within a predictable size, embedding usually keeps reads simple and fast. If the child data is shared, reused, or grows unpredictably, referencing is safer.

This is also a good place to mention that flexible schema can create schema drift if teams do not agree on structure. That is why review practices and validation rules matter.""",
    },
    {
        "kind": "code",
        "tag": "CRUD",
        "title": "MongoDB CRUD Basics",
        "bullets": [
            "`insertOne` and `insertMany` create documents",
            "`find` uses filters and projections to return just the fields you need",
            "`updateOne` and `updateMany` use operators like `$set`, `$inc`, and `$addToSet`",
            "Single-document writes are atomic by default, so document boundaries matter",
            "Always write filters carefully so updates and deletes stay targeted",
        ],
        "code_title": "mongosh examples",
        "code": """db.orders.insertOne({
  customerId: 42,
  status: "pending",
  total: 138.99
});

db.orders.find(
  { status: "pending", total: { $gte: 100 } },
  { customerId: 1, total: 1, _id: 0 }
);

db.orders.updateOne(
  { customerId: 42, status: "pending" },
  { $set: { status: "paid" }, $inc: { retryCount: 1 } }
);""",
        "footer": "Teach students to read the filter first, then the projection or update operator.",
        "notes": """Do not let this become a syntax dump. The educational goal is to show the shape of MongoDB operations, not to memorize every operator.

A reliable reading strategy helps: first identify the target collection, then the filter, then what fields are returned or modified. That mirrors the SQL habit students already have.

Also stress that document boundaries influence correctness. Since single-document writes are atomic, the document structure itself becomes part of the transaction design decision.""",
    },
    {
        "kind": "code",
        "tag": "Aggregation",
        "title": "How To Think About the Aggregation Pipeline",
        "bullets": [
            "Aggregation transforms documents stage by stage instead of returning raw rows only",
            "A common mental model is: filter, reshape, group, sort, project",
            "Ask for the business question first, then pick the stages needed to answer it",
            "Stages like `$unwind` can multiply work, so be cautious with large arrays",
            "Index the early match and sort fields that narrow the pipeline fastest",
        ],
        "code_title": "Revenue by customer",
        "code": """db.orders.aggregate([
  { $match: { status: "shipped" } },
  { $unwind: "$lineItems" },
  {
    $group: {
      _id: "$customerId",
      totalRevenue: {
        $sum: { $multiply: ["$lineItems.qty", "$lineItems.price"] }
      }
    }
  },
  { $sort: { totalRevenue: -1 } },
  { $limit: 5 }
]);""",
        "footer": "Aggregation is easier to teach when every stage answers a visible business question.",
        "notes": """Students coming from SQL often find the pipeline unusual at first. Help them by narrating each stage in plain language. Match shipped orders. Unwind line items. Group by customer. Sum revenue. Sort. Limit.

That narration matters because it keeps the pipeline tied to intent rather than syntax. It also creates a bridge to SQL's mental model of filter, group, aggregate, and order.

Use this slide to remind them that performance is still a DBA concern. Flexibility does not excuse poor query planning or weak indexing.""",
    },
    {
        "kind": "content",
        "tag": "Trade-offs",
        "title": "When MongoDB Helps and When Relational Still Wins",
        "bullets": [
            "MongoDB shines for evolving, nested, product-like, profile-like, or content-heavy data",
            "PostgreSQL still wins for complex ad hoc joins, rigid integrity, and mature relational analytics",
            "Denormalization can speed reads, but it creates synchronization risk when shared facts change",
            "MongoDB supports multi-document transactions, but needing them everywhere can signal a poor fit",
            "Choose the database by workload fit, not by the fashion of the moment",
        ],
        "terms": ["denormalization", "transaction boundary", "join-heavy", "consistency", "workload fit"],
        "callout": "Best tool for the job: the right answer is often PostgreSQL plus MongoDB, not PostgreSQL or MongoDB.",
        "footer": "Trade-off literacy matters more than memorizing a vendor feature list.",
        "notes": """This slide is where you protect students from oversimplified narratives. MongoDB is useful, but it is not the answer to every data problem. The best engineers learn to say why a tool fits and why another tool does not.

Use concrete examples. A catalog with variable product attributes often fits MongoDB well. A financial system with heavy cross-table reporting and strict relational constraints often still fits PostgreSQL better.

If time allows, ask students to name one project idea that is more relational and one that is more document-shaped. That simple exercise reveals whether they are internalizing the trade-offs.""",
    },
    {
        "kind": "content",
        "tag": "Admin",
        "title": "DBA Responsibilities Do Not Disappear in MongoDB",
        "bullets": [
            "Turn rules into validators, unique indexes, role boundaries, and disciplined review processes",
            "Design indexes from actual query patterns instead of indexing every field 'just in case'",
            "Monitor memory, page faults, replication lag, disk growth, slow queries, and connection pressure",
            "Test backup, restore, and failover paths instead of only trusting that backups exist",
            "Secure the cluster with authentication, network controls, TLS, and least privilege",
        ],
        "terms": ["validator", "replication lag", "restore test", "RBAC", "slow query"],
        "callout": (
            "Flexible schema moves enforcement work; it does not eliminate enforcement work."
        ),
        "footer": "Administration shifts shape, but the DBA mindset of correctness and recovery stays the same.",
        "notes": """The course is Database Administration, so keep returning to the operator perspective. Students should not leave with the impression that MongoDB is only about developer convenience.

The right teaching contrast is this: in PostgreSQL, many structural rules are expressed through tables, constraints, and relational design. In MongoDB, some of that enforcement shifts into JSON Schema validators, application review, and query-driven index discipline. The need for governance is still there.

This is also a good place to connect back to earlier weeks. Recovery, least privilege, and observability were not 'Postgres-only' ideas. They are database administration habits that transfer across systems.""",
    },
    {
        "kind": "content",
        "tag": "Pitfalls",
        "title": "Common Misconceptions and Anti-Patterns",
        "bullets": [
            "NoSQL does not mean no schema; uncontrolled schema drift is still a real problem",
            "Giant documents and unbounded arrays can become performance and maintenance hazards",
            "Horizontal scale is architecture work, not a magic setting you flip on later",
            "Using MongoDB for every relationship-heavy workflow can recreate join pain in new forms",
            "If every feature needs broad cross-document transactions, revisit the data model choice",
        ],
        "terms": ["schema drift", "unbounded array", "hotspot", "anti-pattern", "governance"],
        "callout": "Bad modeling in MongoDB does not disappear. It just fails differently than bad SQL modeling.",
        "footer": "Teach students to avoid hype-driven design and to notice where document models can go wrong.",
        "notes": """This slide helps students develop skepticism in the healthy sense. Every database choice creates failure modes.

Explain that MongoDB can encourage fast iteration, but that same flexibility can produce inconsistent field names, mixed document shapes, giant arrays, and painful migrations if teams are careless. Those are design failures, not product bugs.

If students can name these anti-patterns now, they will make better project choices later.""",
    },
    {
        "kind": "content",
        "tag": "Workshop",
        "title": "Classify the Workload Before Naming a Product",
        "bullets": [
            "Shopping catalog with variable product attributes and rich product pages",
            "Session store that needs extremely fast key-based lookups",
            "Fraud detection across accounts, devices, emails, and shared addresses",
            "Sensor telemetry written continuously and queried by recent time windows",
            "For each scenario, choose a model family and name one admin concern",
        ],
        "terms": ["document", "key-value", "graph", "column-family", "SLA"],
        "callout": "Require answers in this form: workload, query shape, consistency need, then database family.",
        "footer": "This is where students prove they can reason about database fit instead of reciting definitions.",
        "notes": """Use this slide as a think-pair-share or short cold-call exercise. It converts the lecture from passive recognition into applied reasoning.

The important thing is not the exact product name students choose. What matters is whether they justify the choice using access pattern, relationship shape, growth pattern, and operational concerns.

Ask for one admin concern with every answer. That keeps the class anchored in DBA thinking instead of drifting into pure application design.""",
    },
    {
        "kind": "content",
        "tag": "Review",
        "title": "Exit Ticket and Next Step",
        "bullets": [
            "In one sentence, explain why NoSQL became prominent again in the 2000s",
            "Name the four major NoSQL families and one fitting workload for each",
            "Choose whether your current project idea is more relational or document-shaped and defend it",
            "Write one embed-versus-reference decision you would make in MongoDB",
            "Use `mongodb_nosql_review_packet.md` for retrieval practice after class",
        ],
        "terms": ["retrieval practice", "project fit", "embed", "reference", "trade-off"],
        "callout": "The fastest way to forget a lecture is to only reread it. Retrieval practice forces understanding.",
        "footer": "Close the session by forcing recall, comparison, and a concrete modeling decision.",
        "notes": """The exit ticket turns the lecture into learning evidence. Students should have to produce concise reasoning, not just nod along during class.

The review packet extends that work after class. Its purpose is to make students retrieve the historical story, the NoSQL families, the MongoDB modeling rules, and the administrator responsibilities without relying only on slide review.

If time is tight, at least collect a short written answer to the project-fit prompt. That answer tells you whether students are actually ready to choose between relational and document models.""",
    },
]


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()

    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(0.42), Inches(3.45), Inches(0.46)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = NAVY_LIGHT
    glow.line.color.rgb = TEAL
    glow.line.width = Pt(1.0)


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(11.85), Inches(0.26))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_notes(slide, notes_text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes_text.strip()


def add_tag(slide, text, width=1.75):
    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.45), Inches(width), Inches(0.38)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_box.line.width = Pt(1.0)

    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL


def add_title_text(slide, text, width=10.0):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), Inches(width), Inches(0.75))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE


def populate_bullets(text_frame, bullets, font_size=19, color=WHITE, bullet_prefix="- "):
    text_frame.word_wrap = True
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"{bullet_prefix}{bullet}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)


def add_terms_panel(slide, terms):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.78), Inches(1.58), Inches(3.62), Inches(4.9)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.3)

    tf = panel.text_frame
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = "Key terms"
    header.font.name = "Aptos"
    header.font.size = Pt(16)
    header.font.bold = True
    header.font.color.rgb = GOLD

    for term in terms:
        p = tf.add_paragraph()
        p.text = term
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(5)


def add_callout(slide, text):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(6.28), Inches(11.52), Inches(0.46)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY_LIGHT
    box.line.color.rgb = ROSE
    box.line.width = Pt(1.2)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.62), Inches(1.9), Inches(0.44)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_box.line.width = Pt(1.2)

    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = spec["tag"]
    tag_p.alignment = PP_ALIGN.CENTER
    tag_p.font.name = "Aptos"
    tag_p.font.size = Pt(14)
    tag_p.font.bold = True
    tag_p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.28), Inches(10.8), Inches(1.45))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = spec["title"]
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.62), Inches(9.8), Inches(1.42))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    p = subtitle_tf.paragraphs[0]
    p.text = spec["subtitle"]
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(3.82), Inches(3.35), Inches(1.96)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.4)

    tf = panel.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = spec["focus_title"]
    p1.font.name = "Aptos"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = GOLD

    for line in spec["focus_lines"]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_content_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, spec["tag"])
    add_title_text(slide, spec["title"])

    bullet_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(7.35), Inches(4.6))
    populate_bullets(bullet_box.text_frame, spec["bullets"], font_size=18)
    add_terms_panel(slide, spec["terms"])

    if spec.get("callout"):
        add_callout(slide, spec["callout"])

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_compare_panel(slide, left, top, width, height, title, bullets):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.2)

    tf = panel.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.space_after = Pt(8)

    for bullet in bullets:
        bp = tf.add_paragraph()
        bp.text = f"- {bullet}"
        bp.font.name = "Aptos"
        bp.font.size = Pt(16)
        bp.font.color.rgb = WHITE
        bp.space_after = Pt(7)


def add_compare_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, spec["tag"], width=1.65)
    add_title_text(slide, spec["title"], width=11.0)

    add_compare_panel(
        slide,
        Inches(0.85),
        Inches(1.82),
        Inches(5.7),
        Inches(4.55),
        spec["left_title"],
        spec["left_bullets"],
    )
    add_compare_panel(
        slide,
        Inches(6.8),
        Inches(1.82),
        Inches(5.6),
        Inches(4.55),
        spec["right_title"],
        spec["right_bullets"],
    )

    if spec.get("callout"):
        add_callout(slide, spec["callout"])

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_code_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, spec["tag"], width=1.6)
    add_title_text(slide, spec["title"], width=10.7)

    bullet_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(4.55), Inches(4.72))
    populate_bullets(bullet_box.text_frame, spec["bullets"], font_size=17)

    code_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.62), Inches(1.72), Inches(6.7), Inches(4.9)
    )
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = NAVY_PANEL
    code_panel.line.color.rgb = TEAL
    code_panel.line.width = Pt(1.2)

    tf = code_panel.text_frame
    tf.word_wrap = True
    tf.clear()

    header = tf.paragraphs[0]
    header.text = spec["code_title"]
    header.font.name = "Aptos"
    header.font.size = Pt(16)
    header.font.bold = True
    header.font.color.rgb = TEAL
    header.space_after = Pt(8)

    body = tf.add_paragraph()
    body.text = spec["code"]
    body.font.name = "Courier New"
    body.font.size = Pt(13)
    body.font.color.rgb = WHITE

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        if spec["kind"] == "title":
            add_title_slide(prs, spec)
        elif spec["kind"] == "content":
            add_content_slide(prs, spec)
        elif spec["kind"] == "compare":
            add_compare_slide(prs, spec)
        elif spec["kind"] == "code":
            add_code_slide(prs, spec)
        else:
            raise ValueError(f"Unsupported slide kind: {spec['kind']}")

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Wrote {OUTPUT}")
