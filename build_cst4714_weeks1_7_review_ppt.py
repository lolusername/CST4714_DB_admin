from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("cst4714_weeks1_7_review_with_script.pptx")

NAVY = RGBColor(6, 16, 31)
NAVY_LIGHT = RGBColor(13, 34, 66)
GOLD = RGBColor(247, 191, 86)
TEAL = RGBColor(123, 224, 214)
WHITE = RGBColor(245, 247, 250)
MUTED = RGBColor(196, 208, 230)


SLIDES = [
    {
        "kind": "title",
        "tag": "Weeks 1-7 Review",
        "title": "CST4714 Review Through Week 7",
        "subtitle": "Relational databases, PostgreSQL, SQL, Supabase, transactions, MVCC, and wraparound.",
        "footer": "Built from the local week folders, labs, PDFs, and textbook topics in this project.",
        "notes": """Today I am going to do a fast review of everything we have covered through week 7 before we open the Jeopardy game. Nothing on this deck is new. The goal is to reconnect the main terms, the main workflows, and the main reasons behind them.

As we go, listen for the same words that appear in the game: database, schema, view, RLS, COMMIT, deadlock, MVCC, autovacuum, wraparound, and correctness. If a term feels familiar but fuzzy, that is okay. The point of this review is to make it sharp again before we compete.

I want you to notice one pattern through the whole deck. We keep returning to the same question: how do we keep data correct, available, and manageable while real users and real changes keep happening?""",
    },
    {
        "kind": "content",
        "tag": "Roadmap",
        "title": "What We Have Covered So Far",
        "bullets": [
            "Week 1: DBA role, DBMS purpose, relational basics, normalization",
            "Week 2: Docker, local Postgres, pgAdmin, and Supabase",
            "Week 3: admin workflows, backups, roles, least privilege, and RLS",
            "Week 4: views, materialized views, identity, constraints, introspection",
            "Week 5: transactions, isolation, blocking, deadlocks, diagnostics",
            "Weeks 6-7: MVCC, vacuum, autovacuum, wraparound, Sentry case study",
        ],
        "terms": ["database", "schema", "migration", "blocking", "vacuum"],
        "footer": "The local week folders are the source of truth for this review.",
        "notes": """This is the map for the review. Week 1 gave us the DBA role, DBMS basics, the relational model, and normalization. Week 2 took PostgreSQL from local tools to Supabase. Week 3 shifted into administration and DevOps habits in a managed Postgres environment. Week 4 focused on schema management: views, materialized views, identity, constraints, and introspection. Week 5 moved into transactions, isolation, locking, and diagnostics. Weeks 6 and 7 pushed deeper into MVCC, vacuum, autovacuum, and the Sentry wraparound case.

So the big story is simple. We started with structure. Then we added tools. Then we added operational discipline. Then we added concurrency and maintenance. By the time we reach week 7, the course is not just asking what a table is. It is asking how a real PostgreSQL system stays correct under load, change, and failure.""",
    },
    {
        "kind": "content",
        "tag": "Week 1",
        "title": "DBMS and DBA Foundations",
        "bullets": [
            "A DBMS helps with consistency, concurrency, recovery, and security",
            "The DBA plans, builds, tunes, troubleshoots, secures, and monitors",
            "A database is the container you connect to",
            "A schema is the namespace inside that database",
            "In a three-tier system: client, app server, database server",
        ],
        "terms": ["DBMS", "DBA", "database", "schema", "app server"],
        "footer": "Weeks 1 and 3 both stress responsibility, not just syntax.",
        "notes": """Here is the first big idea. A DBMS exists because file-based systems are hard to keep consistent, secure, and recoverable. A database management system gives us centralized control over data, better concurrency control, recovery after failures, and stronger security.

The DBA's job is not just to write queries. The DBA plans the system, builds structures, tunes performance, troubleshoots problems, secures access, and monitors for risk. When I say database administrator in this class, I want you to hear both design responsibility and operational responsibility.

Another key distinction is database versus schema. A database is the container you connect to. A schema is the namespace inside that database. If you mix those up, the rest of PostgreSQL gets confusing fast. Also remember the three-tier picture: client on one side, application logic in the middle, and the database server behind it.""",
    },
    {
        "kind": "content",
        "tag": "Weeks 1-2",
        "title": "Relational Model and SQL Basics",
        "bullets": [
            "Relation means table, tuple means row, attribute means column",
            "Domains define what values are allowed in a column",
            "Primary keys identify rows; foreign keys connect tables",
            "Core SQL verbs and clauses: SELECT, WHERE, JOIN, GROUP BY",
            "The class practice dataset uses members, facilities, and bookings",
        ],
        "terms": ["relation", "primary key", "foreign key", "JOIN", "GROUP BY"],
        "footer": "The pgExercises dataset gave us repeated examples for joins and summaries.",
        "notes": """When you hear relational model, think tables with rules. A relation is a table. A tuple is a row. An attribute is a column. A domain is the legal set of values that a column is allowed to hold. Those terms matter because they give us precise language for talking about design.

Keys are just as important. A primary key uniquely identifies a row. A foreign key connects a child table to a parent table and helps enforce referential integrity. That is how the database knows that relationships are valid, not just the application.

On the SQL side, the core language we kept returning to was SELECT, WHERE, JOIN, and GROUP BY. We practiced that thinking with the Club Database, especially the members, facilities, and bookings tables. Also keep normalization in mind: the reason we normalize is to reduce redundancy and avoid update problems later.""",
    },
    {
        "kind": "content",
        "tag": "Week 2",
        "title": "From Local Postgres to Supabase",
        "bullets": [
            "Docker gave us repeatable local Postgres environments",
            "Image, container, volume, and network are the key Docker words",
            "We connected with both psql and pgAdmin",
            "Database and schema are not the same thing",
            "Supabase is still PostgreSQL, but with managed infrastructure",
        ],
        "terms": ["Docker", "container", "volume", "pgAdmin", "Supabase"],
        "footer": "Local tools changed, but PostgreSQL language and logic stayed the same.",
        "notes": """Week 2 was about portability. Docker let us run the same PostgreSQL setup on different machines with fewer environment surprises. The key Docker words were image, container, volume, and network. The one students most often forget is volume, and that matters because the volume is what keeps data alive when containers come and go.

We also connected to Postgres in two ways: psql in the terminal and pgAdmin in a GUI. That mattered because good DBAs should be comfortable in both styles.

This week also reinforced the database versus schema distinction. And then we moved to Supabase. The big lesson there is that local does not mean one database and cloud means another database product. Supabase is still PostgreSQL. The engine is the same. What changes is the platform layer, the tooling, and who manages which responsibilities.""",
    },
    {
        "kind": "content",
        "tag": "Week 3",
        "title": "Supabase Operations and Shared Responsibility",
        "bullets": [
            "Managed Postgres does not remove responsibility for correctness",
            "Use logs, reports, backups, settings, Auth, and RLS intentionally",
            "Plan, migrate, verify, observe, and recover for every change",
            "Think restore-first, not just backup-first",
            "Use least privilege for roles, grants, and exposed access paths",
        ],
        "terms": ["shared responsibility", "migration-first", "PITR", "least privilege", "RLS"],
        "footer": "Week 3 moved us from user mindset to operator mindset.",
        "notes": """Week 3 changed the mindset from student user to database operator. The central message was simple: managed Postgres means fewer server chores, not less accountability. Supabase may handle infrastructure layers, but you still own schema quality, query quality, index quality, migration safety, privileges, RLS logic, and recovery verification.

That is why the course kept returning to the same workflow: plan the change, apply the change through a migration, verify the change with SQL and evidence, observe for regressions, and recover if needed. This is migration-first and verification-first discipline.

We also talked about least privilege and RLS. Least privilege means give only the permissions needed, not broad access by default. RLS means row access should be filtered intentionally. One specific warning from class is that the service role bypasses RLS, so it must not be treated casually. These are not optional extras. They are part of safe operations.""",
    },
    {
        "kind": "content",
        "tag": "Week 4",
        "title": "Schema Management as Policy",
        "bullets": [
            "A view is a saved query that acts like a virtual table",
            "A materialized view stores results and must be refreshed",
            "IDENTITY is the modern auto-generated ID pattern",
            "Constraints are data integrity contracts, not decorations",
            "Use information_schema and pg_catalog to inspect live reality",
        ],
        "terms": ["view", "materialized view", "IDENTITY", "CHECK", "information_schema"],
        "footer": "Week 4 treated structure as something the database actively enforces.",
        "notes": """Week 4 treated structure as policy. A view is not just a convenience. It is a saved query that can give consumers a stable interface. A materialized view is different because it stores results physically and has to be refreshed. The tradeoff is freshness versus speed.

For key generation, we emphasized that IDENTITY is the modern, SQL-standard choice for new Postgres tables. The reason is not fashion. It is clarity and better alignment with modern schema practice.

Constraints were one of the biggest ideas of the week. A CHECK constraint, a UNIQUE constraint, a primary key, and a foreign key are all ways the database enforces business truth at write time. They turn silent bad data into explicit errors.

Finally, introspection matters because assumptions are not enough. information_schema is the standards-oriented place to inspect metadata, while pg_catalog gives PostgreSQL-specific detail when you need deeper answers.""",
    },
    {
        "kind": "content",
        "tag": "Week 5",
        "title": "Transactions and Isolation",
        "bullets": [
            "Autocommit treats each statement as its own transaction",
            "Explicit transactions group multiple steps into one correctness unit",
            "COMMIT saves work; ROLLBACK undoes uncommitted work",
            "READ COMMITTED is the default isolation level in PostgreSQL",
            "Keep transactions short, especially in pooled environments",
        ],
        "terms": ["BEGIN", "COMMIT", "ROLLBACK", "READ COMMITTED", "SERIALIZABLE"],
        "footer": "Week 5 added time and overlapping work to the model.",
        "notes": """Week 5 added time to the picture. Up to this point, a lot of design thinking can feel static. Transactions force us to ask what happens when multiple valid actions overlap.

The first idea is transaction boundary. Autocommit means each statement is its own transaction. Explicit transactions let us say these several steps either all happen together or none of them count. COMMIT makes the work permanent. ROLLBACK throws away the uncommitted work.

The second idea is isolation. PostgreSQL defaults to READ COMMITTED, which means each statement sees the latest committed state at the start of that statement. That is fine for many workloads, but not all. If a workflow has stronger correctness needs, we may need a stronger pattern or a stronger isolation level.

The course kept saying the same thing here: a transaction is not just SQL grouping. It is a correctness boundary. And short transactions are safer than long ones.""",
    },
    {
        "kind": "content",
        "tag": "Week 5",
        "title": "Blocking, Deadlocks, and Diagnostics",
        "bullets": [
            "Locks are normal when conflicting writes touch the same rows",
            "Blocking means one transaction is waiting on another",
            "A deadlock is a cycle of waiting transactions",
            "Use pg_stat_activity and wait information to inspect the problem",
            "Mitigation restores service; prevention fixes the underlying pattern",
        ],
        "terms": ["blocking", "deadlock", "pg_stat_activity", "wait_event", "incident response"],
        "footer": "The lab focused on evidence, not guessing.",
        "notes": """Locks are normal. Blocking is also normal in the sense that it can happen whenever two transactions want incompatible access to the same data. The problem is not that locks exist. The problem is long lock hold times, bad ordering, or transactions staying open longer than they should.

In lab, we used two sessions to create blocking on purpose and then inspected it with pg_stat_activity. That view shows active sessions, state, wait information, and the running query text. It is one of the first places to look when users say a query is hanging.

A deadlock is more specific. It is a cycle where transaction A waits on transaction B while transaction B waits on transaction A. PostgreSQL resolves that by aborting one transaction. Operationally, that means mitigation and root cause are different. Mitigation gets service moving again. Root cause changes the design so the same pattern stops repeating.""",
    },
    {
        "kind": "content",
        "tag": "Week 6",
        "title": "MVCC, Vacuum, and Autovacuum",
        "bullets": [
            "MVCC lets readers and writers overlap more safely",
            "Updates create new row versions instead of overwriting in place",
            "Old row versions become dead tuples",
            "VACUUM reclaims reusable space; ANALYZE refreshes planner stats",
            "Long transactions and neglected cleanup create operational risk",
        ],
        "terms": ["MVCC", "dead tuples", "vacuum", "autovacuum", "analyze"],
        "footer": "Maintenance affects both performance and long-term safety.",
        "notes": """MVCC stands for Multi-Version Concurrency Control. The simple mental model is that PostgreSQL keeps row versions so that readers usually do not block writers and writers do not usually block readers the way students first imagine.

When a row is updated, PostgreSQL creates a new version rather than simply erasing the old one in place. The old version can become a dead tuple after it is no longer needed for visibility. That is where vacuum comes in. Vacuum reclaims reusable space and helps keep tables healthy. Analyze updates the statistics the planner uses to estimate row counts and choose plans.

Autovacuum is the background process that automates much of this work. The class kept stressing that vacuum is not only a performance detail. It is part of correctness and long-term health. Also remember the risk of long transactions. They can delay cleanup because older snapshots may still need older row versions to remain visible.""",
    },
    {
        "kind": "content",
        "tag": "Week 7",
        "title": "Wraparound and the Sentry Incident",
        "bullets": [
            "Transaction IDs are finite, so PostgreSQL must prevent dangerous wraparound",
            "If wraparound risk gets too high, PostgreSQL may stop writes",
            "That is a correctness protection, not just a performance slowdown",
            "Sentry traded some data in one noncritical table for faster recovery",
            "Reading URL: blog.sentry.io/transaction-id-wraparound-in-postgres/",
            "Prevention means monitoring, vacuum discipline, and operational readiness",
        ],
        "terms": ["XID", "wraparound", "correctness", "availability", "TRUNCATE"],
        "footer": "Week 7 turned maintenance into a real production story.",
        "notes": """The Sentry reading matters because it shows why PostgreSQL maintenance is not just housekeeping. Transaction IDs are finite. If old rows are not cleaned up properly and the system approaches wraparound danger, PostgreSQL may stop accepting writes. It does that to protect correctness.

The exact article URL for this activity is https://blog.sentry.io/transaction-id-wraparound-in-postgres/

That is the key phrase: protect correctness. This is not only about the database running slowly. This is about preserving valid visibility rules so the database does not lie about what rows should be seen.

In the Sentry incident, the team chose to truncate one large noncritical table so the service could recover faster. That was a tradeoff between availability and preserving every piece of data in that one table. The lesson is not that truncation is always right. The lesson is that operational decisions under pressure are about priorities, evidence, and understanding what the database is protecting.""",
    },
    {
        "kind": "content",
        "tag": "Game Prep",
        "title": "What To Listen For In Jeopardy",
        "bullets": [
            "Direct pairs: view to saved query, RLS to row access, pg_stat_activity to active sessions",
            "Remember the repeated examples: members, facilities, bookings",
            "If a clue sounds operational, ask what risk or responsibility it points to",
            "If a clue sounds about concurrency, ask what protects correctness under overlap",
            "After this review, open cst4714_weeks1_7_jeopardy.html",
        ],
        "terms": [
            "DBA Foundations",
            "SQL + Relational",
            "Supabase + Ops",
            "Schema Management",
            "Transactions + Locking",
            "MVCC + Vacuum",
        ],
        "footer": "Big theme for the whole board: correctness under change, access, and concurrency.",
        "notes": """When we open the game, do not overcomplicate the clues. Many of them are direct recall. If you hear saved query that looks like a table, think view. If you hear row access in Supabase, think RLS. If you hear active sessions and waiting queries, think pg_stat_activity. If you hear automatic cleanup, think autovacuum.

Also remember the repeated examples from class. The members, facilities, and bookings tables showed up again and again because they gave us a concrete relational schema for joins, views, summaries, transactions, and incident drills.

Most importantly, if you get stuck on a clue, ask what problem the feature solves. Does it solve data modeling, access control, concurrency, maintenance, or recovery? That question usually points you to the answer. The bigger theme behind the whole board is correctness under change, access, and concurrency.""",
    },
]


def set_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.fill.background()

    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(0.45), Inches(3.3), Inches(0.42)
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = NAVY_LIGHT
    glow.line.color.rgb = TEAL
    glow.line.width = Pt(1.0)


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.75), Inches(7.02), Inches(11.8), Inches(0.25))
    tf = footer.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_notes(slide, notes_text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = notes_text.strip()


def add_title_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.62), Inches(1.9), Inches(0.44)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_box.line.width = Pt(1.2)
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = spec["tag"]
    tag_p.alignment = PP_ALIGN.CENTER
    tag_p.font.name = "Aptos"
    tag_p.font.size = Pt(14)
    tag_p.font.bold = True
    tag_p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(10.9), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = spec["title"]
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.72), Inches(9.7), Inches(1.4))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    p = subtitle_tf.paragraphs[0]
    p.text = spec["subtitle"]
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.25), Inches(3.9), Inches(3.15), Inches(1.6)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = NAVY_LIGHT
    right_panel.line.color.rgb = GOLD
    right_panel.line.width = Pt(1.4)

    rp_tf = right_panel.text_frame
    rp_tf.word_wrap = True
    p1 = rp_tf.paragraphs[0]
    p1.text = "Review focus"
    p1.font.name = "Aptos"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = GOLD

    p2 = rp_tf.add_paragraph()
    p2.text = "Structure, operations, concurrency, and correctness"
    p2.font.name = "Aptos"
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def add_content_slide(prs, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(0.45), Inches(1.65), Inches(0.38)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = NAVY_LIGHT
    tag_box.line.color.rgb = TEAL
    tag_tf = tag_box.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = spec["tag"]
    tag_p.alignment = PP_ALIGN.CENTER
    tag_p.font.name = "Aptos"
    tag_p.font.size = Pt(12)
    tag_p.font.bold = True
    tag_p.font.color.rgb = TEAL

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.92), Inches(8.3), Inches(0.75))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = spec["title"]
    p.font.name = "Aptos Display"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    bullet_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(7.35), Inches(4.95))
    bullet_tf = bullet_box.text_frame
    bullet_tf.word_wrap = True
    bullet_tf.clear()

    for index, bullet in enumerate(spec["bullets"]):
        paragraph = bullet_tf.paragraphs[0] if index == 0 else bullet_tf.add_paragraph()
        paragraph.text = f"- {bullet}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(10)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(1.55), Inches(3.7), Inches(4.95)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY_LIGHT
    panel.line.color.rgb = GOLD
    panel.line.width = Pt(1.4)

    panel_tf = panel.text_frame
    panel_tf.word_wrap = True
    header = panel_tf.paragraphs[0]
    header.text = "Key terms"
    header.font.name = "Aptos"
    header.font.size = Pt(16)
    header.font.bold = True
    header.font.color.rgb = GOLD

    for term in spec["terms"]:
        p = panel_tf.add_paragraph()
        p.text = term
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)

    add_footer(slide, spec["footer"])
    add_notes(slide, spec["notes"])


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        if spec["kind"] == "title":
            add_title_slide(prs, spec)
        else:
            add_content_slide(prs, spec)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Wrote {OUTPUT}")
